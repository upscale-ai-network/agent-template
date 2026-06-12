"""Low-res PNG previews from built PPTX — proportional downscale, no re-layout."""

from __future__ import annotations

import io
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]

FOOTER_TOP = 4_800_000
DEFAULT_WIDTH = 960
GRID_GAP = 8
EMU_PER_INCH = 914400


def _shape_rgb(shape) -> Tuple[int, int, int]:
    if not shape.has_text_frame:
        return (5, 24, 48)
    try:
        para = shape.text_frame.paragraphs[0]
        if para.runs and para.runs[0].font.color.rgb is not None:
            rgb: RGBColor = para.runs[0].font.color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except (AttributeError, IndexError):
        pass
    if shape.top < 700_000:
        return (5, 24, 48)
    return (51, 51, 51)


def _paragraph_size_pt(paragraph) -> float:
    for run in paragraph.runs:
        if run.text.strip() and run.font.size is not None:
            return float(run.font.size.pt)
    return 11.0


def _font_px(pt: float, *, width: int, slide_width: int) -> int:
    slide_px = slide_width / EMU_PER_INCH * 96.0
    return max(6, int(round(pt * (96.0 / 72.0) * width / slide_px)))


def _font_for_pt(pt: float, *, width: int, slide_width: int):
    px = _font_px(pt, width=width, slide_width=slide_width)
    for path in (
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ):
        try:
            return ImageFont.truetype(path, px)
        except OSError:
            continue
    return ImageFont.load_default()


def _render_text_shape(
    draw,
    shape,
    *,
    width: int,
    slide_width: int,
    scale: float,
) -> None:
    if not shape.has_text_frame:
        return
    x = int(shape.left * scale)
    y = int(shape.top * scale)
    max_h = max(1, int(shape.height * scale))
    rgb = _shape_rgb(shape)
    line_gap = max(1, int(round(2 * scale * EMU_PER_INCH / 12700)))
    cursor_y = y
    for para in shape.text_frame.paragraphs:
        text = para.text
        if not text.strip():
            cursor_y += line_gap
            continue
        pt = _paragraph_size_pt(para)
        font = _font_for_pt(pt, width=width, slide_width=slide_width)
        draw.text((x, cursor_y), text, fill=rgb, font=font)
        bbox = draw.textbbox((x, cursor_y), text, font=font)
        cursor_y = bbox[3] + line_gap
        if cursor_y > y + max_h + line_gap * 4:
            break


def _render_picture(img, shape, *, scale: float) -> None:
    blob = shape.image.blob
    pic = Image.open(io.BytesIO(blob)).convert("RGBA")
    w = max(1, int(shape.width * scale))
    h = max(1, int(shape.height * scale))
    pic = pic.resize((w, h), Image.Resampling.LANCZOS)
    x = int(shape.left * scale)
    y = int(shape.top * scale)
    img.paste(pic, (x, y), pic)


def _render_slide_faithful(slide, *, width: int, slide_width: int, slide_height: int):
    scale = width / float(slide_width)
    height = max(1, int(round(slide_height * scale)))
    img = Image.new("RGB", (width, height), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    shapes = list(slide.shapes)
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            _render_picture(img, shape, scale=scale)
    for shape in shapes:
        if shape.has_text_frame and shape.top <= FOOTER_TOP:
            _render_text_shape(
                draw, shape, width=width, slide_width=slide_width, scale=scale
            )
    return img


def _build_grid(slide_paths: List[Path], grid_path: Path, *, width: int) -> Path:
    if not slide_paths:
        return grid_path
    images = [Image.open(p) for p in slide_paths]
    thumb_h = max(im.height for im in images)
    cols = min(4, len(images))
    rows = (len(images) + cols - 1) // cols
    grid_w = cols * width + (cols + 1) * GRID_GAP
    grid_h = rows * thumb_h + (rows + 1) * GRID_GAP
    grid = Image.new("RGB", (grid_w, grid_h), (240, 244, 248))
    for idx, im in enumerate(images):
        r, c = divmod(idx, cols)
        x = GRID_GAP + c * (width + GRID_GAP)
        y = GRID_GAP + r * (thumb_h + GRID_GAP)
        grid.paste(im, (x, y))
    grid.save(grid_path, "PNG", optimize=True)
    for im in images:
        im.close()
    return grid_path


def export_pptx_previews(
    pptx_path: Path,
    out_dir: Optional[Path] = None,
    *,
    width: int = DEFAULT_WIDTH,
) -> List[Path]:
    """Write slide-NN.png + {stem}-grid.png — scaled raster of PPTX layout."""
    pptx_path = Path(pptx_path)
    if not pptx_path.is_file():
        raise FileNotFoundError(pptx_path)

    out_dir = Path(out_dir or pptx_path.parent / "preview")
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    sw, sh = prs.slide_width, prs.slide_height
    slide_paths: List[Path] = []

    for i, slide in enumerate(prs.slides, start=1):
        img = _render_slide_faithful(slide, width=width, slide_width=sw, slide_height=sh)
        out = out_dir / f"slide-{i:02d}.png"
        img.save(out, "PNG", optimize=True)
        slide_paths.append(out)

    grid_path = out_dir / f"{pptx_path.stem}-grid.png"
    _build_grid(slide_paths, grid_path, width=width)
    return slide_paths + [grid_path]


def main() -> int:
    import argparse

    parser = argparse.ArgumentParser(description="Export low-res PNG previews from PPTX.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("-o", "--out-dir", type=Path, default=None)
    parser.add_argument("-W", "--width", type=int, default=DEFAULT_WIDTH)
    args = parser.parse_args()
    paths = export_pptx_previews(args.pptx, args.out_dir, width=args.width)
    for p in paths:
        label = p.relative_to(ROOT) if p.is_relative_to(ROOT) else p
        print(f"OK: {label}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
