#!/usr/bin/env python3
"""
Build DT100 A3/B6 PPTX from bugatti-qos-architecture.md and bugatti-qos-ccc.md.
"""

import shutil
import sys
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
DT100 = ROOT / "dt100"
DT122 = ROOT / "dt122"
sys.path.insert(0, str(ROOT / "scripts"))

from deck_from_md import A3_MD, B6_MD, load_b6_md, load_deck_md  # noqa: E402
from deck_validate import (  # noqa: E402
    DeckBuildError,
    fail_on_errors,
    validate_a3_build,
    validate_a3_diagram_pngs,
    validate_b6_build,
    validate_b6_pptx_no_speaker_notes,
    validate_built_pptx,
)
from pptx_util import (  # noqa: E402
    assert_pptx_valid,
    fill_content_diagram_slide,
    fill_cover_slide,
    save_presentation,
    scrub_speaker_notes,
    trim_to_slides,
)

PIPELINE_IMG = ROOT / "assets" / "logical-pipeline-boss-slide.png"
A3_DIAGRAMS = ROOT / "assets" / "diagrams" / "a3"
B6_DIAGRAMS = ROOT / "assets" / "diagrams" / "b6"

# Committed company template (git). Optional local seed: ~/Downloads/Mirror-*.pptx
STYLE_DOWNLOAD = Path.home() / "Downloads" / "Mirror-Sflow-Bugatti-ASIC-CCC.pptx"
STYLE_REF = ROOT / "assets/templates/upscale-ccc-style-reference.pptx"  # future: JFrog artifact pin

IDX_COVER = 0
IDX_CONTENT = 2


def ensure_style_reference() -> Path:
    if STYLE_REF.exists():
        return STYLE_REF
    if STYLE_DOWNLOAD.exists():
        STYLE_REF.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(STYLE_DOWNLOAD, STYLE_REF)
        return STYLE_REF
    raise FileNotFoundError(
        f"Company style deck not found. Download CCC PPTX to {STYLE_DOWNLOAD} "
        f"or copy to {STYLE_REF}"
    )


class StyledDeck:
    def __init__(self, out_path: Path, num_content_slides: int):
        ref = ensure_style_reference()
        shutil.copy2(ref, out_path)
        self.path = out_path
        self.prs = Presentation(str(out_path))
        keep = [IDX_COVER] + list(range(IDX_CONTENT, IDX_CONTENT + num_content_slides))
        if keep[-1] >= len(self.prs.slides):
            raise ValueError(f"CCC deck needs {keep[-1] + 1} slides, has {len(self.prs.slides)}")
        trim_to_slides(self.prs, keep)
        self._slide_i = 1

    def add_cover(self, headline, subline, meta, tag, notes=None):
        fill_cover_slide(self.prs.slides[0], headline, subline, meta, tag)
        if notes:
            self.prs.slides[0].notes_slide.notes_text_frame.text = notes

    def add_content(
        self,
        title: Optional[str],
        bullets,
        subtitle=None,
        lead=None,
        notes=None,
        diagram: Optional[Path] = None,
        title_lines: Optional[List[str]] = None,
        caption: Optional[str] = None,
        layout: Optional[str] = None,
    ):
        slide = self.prs.slides[self._slide_i]
        self._slide_i += 1
        if diagram is not None:
            if not diagram.is_file():
                raise FileNotFoundError(f"Slide {title!r} requires diagram PNG: {diagram}")
            fill_content_diagram_slide(
                slide,
                title,
                diagram,
                title_lines=title_lines,
                subtitle_line=subtitle,
                lead_line=lead,
                caption_line=caption,
                bullet_lines=bullets or None,
                slide_width=self.prs.slide_width,
                layout=layout,
            )
        else:
            from pptx_util import fill_content_slide

            fill_content_slide(
                slide, title, bullets, subtitle=subtitle, lead=lead, title_lines=title_lines
            )
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return slide

    def add_image_slide(
        self,
        title: Optional[str],
        image_path: Path,
        *,
        subtitle=None,
        lead=None,
        caption: Optional[str] = None,
        bullets=None,
        notes=None,
    ):
        if not image_path.is_file():
            raise FileNotFoundError(f"Slide {title!r} requires image: {image_path}")
        slide = self.add_content(
            title,
            bullets or [],
            subtitle=subtitle,
            lead=lead,
            caption=caption,
            diagram=image_path,
            notes=notes,
        )
        return slide

    def save(self):
        save_presentation(self.prs, self.path)
        assert_pptx_valid(self.path)
        print(f"OK: {self.path.name}")
        return self.path


def ensure_a3_diagrams(doc, diagram_dir: Path | None = None) -> None:
    from render_a3_diagrams import render_all_a3_diagrams

    stems = render_all_a3_diagrams(doc, diagram_dir=diagram_dir)
    print(f"Rendered {len(stems)} A3 diagrams from {A3_MD.name} (PyMuPDF)")
    for stem in stems:
        p = _a3_png(stem, diagram_dir)
        label = p.relative_to(ROOT) if p.is_relative_to(ROOT) else p
        print(f"OK: {label}")
    fail_on_errors(validate_a3_diagram_pngs(doc, diagram_dir=diagram_dir))


def _a3_png(name: str, diagram_dir: Path | None = None) -> Path:
    return (diagram_dir or A3_DIAGRAMS) / f"{name}.png"


def _b6_png(name: str, diagram_dir: Path | None = None) -> Path:
    return (diagram_dir or B6_DIAGRAMS) / f"{name}.png"


def ensure_b6_diagrams(doc, diagram_dir: Path | None = None) -> None:
    from render_b6_diagrams import render_all_b6_diagrams

    override = diagram_dir
    if override is not None:
        override.mkdir(parents=True, exist_ok=True)
        import render_b6_diagrams as r6

        prev = r6.B6_DIR
        r6.B6_DIR = override
        try:
            stems = render_all_b6_diagrams(doc)
        finally:
            r6.B6_DIR = prev
    else:
        stems = render_all_b6_diagrams(doc)

    base = override or B6_DIAGRAMS
    print(f"Rendered {len(stems)} B6 diagrams from {B6_MD.name}")
    for stem in stems:
        png = base / f"{stem}.png"
        try:
            label = png.relative_to(ROOT)
        except ValueError:
            label = png
        print(f"OK: {label}")
    from deck_validate import validate_b6_diagram_pngs

    fail_on_errors(validate_b6_diagram_pngs(doc, diagram_dir=base))


def _join_notes(*parts: str) -> Optional[str]:
    text = "\n\n".join(p.strip() for p in parts if p and p.strip())
    return text or None


def build_a3(
    out_path: Path | None = None, diagram_dir: Path | None = None
) -> Path:
    doc = load_deck_md(A3_MD, a3_cover_fields=True)
    fail_on_errors(validate_a3_build(doc))
    ensure_a3_diagrams(doc, diagram_dir=diagram_dir)
    out = out_path or (DT100 / "bugatti-qos-architecture.pptx")
    deck = StyledDeck(out, num_content_slides=len(doc.slides))

    cov = doc.cover
    deck.add_cover(
        cov.left_title,
        "\n".join(cov.right_lines),
        cov.meta or cov.left_subtitle,
        cov.tag,
        notes=cov.notes or None,
    )

    before = doc.extra_notes.get("before-1", "")
    into_pipeline = doc.extra_notes.get("into-pipeline", "")
    after_pipeline = doc.extra_notes.get("after-pipeline", "")

    for s in doc.ordered_slides():
        notes = s.notes
        if s.number == 1:
            notes = _join_notes(before, notes, into_pipeline)
        elif s.number == 3:
            notes = _join_notes(notes, after_pipeline)
        kwargs = dict(
            title=s.title,
            bullets=[],
            subtitle=s.subtitle or None,
            notes=notes,
        )
        if s.diagram:
            kwargs["diagram"] = _a3_png(s.diagram, diagram_dir)
        deck.add_content(**kwargs)

    return deck.save()


def build_b6(
    out_path: Path | None = None,
    diagram_dir: Path | None = None,
) -> Path:
    doc = load_b6_md()
    fail_on_errors(validate_b6_build(doc))
    ensure_b6_diagrams(doc, diagram_dir=diagram_dir)
    slides = doc.ordered_slides()
    out = out_path or (DT122 / "bugatti-qos-ccc.pptx")
    deck = StyledDeck(out, num_content_slides=len(slides))

    cov = doc.cover
    deck.add_cover(cov.title, cov.subtitle, cov.meta, cov.tag)

    for s in slides:
        if s.diagram:
            deck.add_content(
                s.title,
                s.bullets,
                subtitle=s.subtitle or None,
                lead=s.lead or None,
                caption=s.caption or None,
                diagram=_b6_png(s.diagram, diagram_dir),
                layout=s.layout or None,
            )
        elif s.image:
            img = PIPELINE_IMG if s.image == "logical-pipeline-boss-slide.png" else ROOT / "assets" / s.image
            deck.add_image_slide(
                s.title or None,
                img,
                subtitle=s.subtitle or None,
                lead=s.lead or None,
                caption=s.caption or None,
                bullets=s.bullets,
            )
        else:
            deck.add_content(
                s.title,
                s.bullets,
                subtitle=s.subtitle or None,
                lead=getattr(s, "lead", None) or None,
            )

    scrub_speaker_notes(deck.prs)
    return deck.save()


def main():
    import argparse

    parser = argparse.ArgumentParser(description="Build DT100 deck PPTX from markdown.")
    parser.add_argument("--a3-only", action="store_true", help="Build draft QoS slides only")
    args = parser.parse_args()

    try:
        ref = ensure_style_reference()
        print(f"Style reference: {ref}")

        a3_doc = load_deck_md(A3_MD, a3_cover_fields=True)
        fail_on_errors(validate_a3_build(a3_doc))
        if not args.a3_only:
            b6_doc = load_b6_md()
            fail_on_errors(validate_b6_build(b6_doc))
        print("Pre-build checks: OK")

        a3 = build_a3()
        a3_slides = 1 + len(a3_doc.slides)
        checks = validate_built_pptx(a3, expected_slides=a3_slides)
        print(f"Wrote {a3} ({a3_slides} slides, company chrome)")

        if not args.a3_only:
            b6 = build_b6()
            b6_slides = 1 + len(b6_doc.ordered_slides())
            checks += validate_built_pptx(b6, expected_slides=b6_slides)
            checks += validate_b6_pptx_no_speaker_notes(b6)
            print(f"Wrote {b6} ({b6_slides} slides, company chrome)")

        fail_on_errors(checks)
        print("Post-build checks: OK")

        from pptx_preview import export_pptx_previews

        for deck_path in (a3,) if args.a3_only else (a3, b6):
            previews = export_pptx_previews(deck_path)
            print(f"Preview: {deck_path.parent / 'preview'} ({len(previews)} PNGs)")
            for p in previews:
                label = p.relative_to(ROOT) if p.is_relative_to(ROOT) else p
                print(f"  OK: {label}")
    except DeckBuildError as exc:
        print(exc, file=sys.stderr)
        raise SystemExit(1) from exc


if __name__ == "__main__":
    main()
