"""Workflow regression kit — build arbitrary md + diagram dir → pptx (canary / tests)."""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path
from typing import Callable, List, Optional

ROOT = Path(__file__).resolve().parents[1]
FIXTURES = ROOT / "tests" / "fixtures"
CANARY_MD = FIXTURES / "workflow-canary.md"
CANARY_DIAGRAMS = FIXTURES / "diagrams" / "canary"
CANARY_CONFIG = CANARY_DIAGRAMS / "mermaid-config.json"
CANARY_OUT = ROOT / "tests" / "output" / "workflow-canary.pptx"

sys.path.insert(0, str(ROOT / "scripts"))

from deck_from_md import DeckDocument, load_deck_md  # noqa: E402
from deck_render import render_stems_from_dir  # noqa: E402
from deck_validate import fail_on_errors, validate_png  # noqa: E402


def _load_build_module():
    path = ROOT / "scripts" / "build-dt100-decks.py"
    spec = importlib.util.spec_from_file_location("build_dt100_decks", path)
    mod = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(mod)
    return mod


def diagram_stems(doc: DeckDocument) -> List[str]:
    return [s.diagram for s in doc.ordered_slides() if s.diagram]


def render_doc_diagrams(
    doc: DeckDocument,
    diagram_dir: Path,
    *,
    config_path: Optional[Path] = None,
    extra_renderers: Optional[dict] = None,
) -> List[str]:
    """Render all diagrams declared in doc. extra_renderers: stem → callable."""
    extra_renderers = extra_renderers or {}
    rendered: List[str] = []
    for stem in diagram_stems(doc):
        if stem in extra_renderers:
            extra_renderers[stem]()
        else:
            render_stems_from_dir(diagram_dir, [stem], config_path=config_path)
        fail_on_errors(
            validate_png(diagram_dir / f"{stem}.png", min_width=150, min_height=50)
        )
        rendered.append(stem)
    return rendered


def png_for_stem(diagram_dir: Path, stem: str) -> Path:
    return diagram_dir / f"{stem}.png"


def build_deck_from_md(
    md_path: Path,
    out_pptx: Path,
    diagram_dir: Path,
    *,
    config_path: Optional[Path] = None,
    render: bool = True,
    extra_renderers: Optional[dict] = None,
) -> DeckDocument:
    build = _load_build_module()
    doc = load_deck_md(md_path, stop_at_reference=True, a3_cover_fields=False)
    if render:
        render_doc_diagrams(
            doc, diagram_dir, config_path=config_path, extra_renderers=extra_renderers
        )

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    slides = doc.ordered_slides()
    deck = build.StyledDeck(out_pptx, num_content_slides=len(slides))
    cov = doc.cover
    deck.add_cover(cov.title, cov.subtitle, cov.meta, cov.tag)

    for s in slides:
        if s.diagram:
            deck.add_content(
                s.title,
                s.bullets,
                subtitle=s.subtitle or None,
                lead=s.lead or s.caption or None,
                diagram=png_for_stem(diagram_dir, s.diagram),
            )
        else:
            deck.add_content(
                s.title,
                s.bullets,
                subtitle=s.subtitle or None,
                lead=s.lead or None,
            )
    deck.save()
    return doc


def build_production_b6_isolated(tmp_dir: Path) -> Path:
    """Full md→diagram→pptx for B6 into tmp_dir; never writes dt122/ or assets/diagrams/b6/."""
    build = _load_build_module()
    diagram_dir = tmp_dir / "diagrams" / "b6"
    out_pptx = tmp_dir / "bugatti-qos-ccc.pptx"
    return build.build_b6(out_path=out_pptx, diagram_dir=diagram_dir)


def build_production_a3_isolated(tmp_dir: Path) -> Path:
    """Full md→diagram→pptx for A3 into tmp_dir; never writes dt100/ or assets/diagrams/a3/."""
    build = _load_build_module()
    diagram_dir = tmp_dir / "diagrams" / "a3"
    out_pptx = tmp_dir / "bugatti-qos-architecture.pptx"
    return build.build_a3(out_path=out_pptx, diagram_dir=diagram_dir)


def build_canary(
    md_path: Path = CANARY_MD,
    diagram_dir: Path = CANARY_DIAGRAMS,
    out_pptx: Path = CANARY_OUT,
) -> DeckDocument:
    return build_deck_from_md(
        md_path, out_pptx, diagram_dir, config_path=CANARY_CONFIG
    )


def extract_visible_text(pptx_path: Path) -> List[str]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    chunks: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    chunks.append(text)
    return chunks


def slide_count(pptx_path: Path) -> int:
    from pptx import Presentation

    return len(Presentation(str(pptx_path)).slides)
