#!/usr/bin/env python3
"""
Build Gluon company slide template: 2 seed slides with full Upscale chrome.
Uses copy-and-trim (valid PPTX) — not clone-into-empty-deck.

Usage:
  python3 scripts/build-company-template.py
"""

import shutil
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
STYLE_DOWNLOAD = Path.home() / "Downloads" / "Mirror-Sflow-Bugatti-ASIC-CCC.pptx"
STYLE_REF = ROOT / "assets/templates/upscale-ccc-style-reference.pptx"
OUT = ROOT / "assets/templates/upscale-company-template.pptx"

IDX_COVER = 0
IDX_CONTENT = 2

sys.path.insert(0, str(ROOT / "scripts"))
from pptx_util import (  # noqa: E402
    apply_content_colors,
    apply_cover_colors,
    assert_pptx_valid,
    fill_content_slide,
    fill_cover_slide,
    save_presentation,
    trim_to_slides,
)


def ensure_style_reference() -> Path:
    if STYLE_REF.exists():
        return STYLE_REF
    if STYLE_DOWNLOAD.exists():
        STYLE_REF.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(STYLE_DOWNLOAD, STYLE_REF)
        return STYLE_REF
    raise FileNotFoundError(f"CCC style deck required: {STYLE_DOWNLOAD}")


def main() -> int:
    ref = ensure_style_reference()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(ref, OUT)

    prs = Presentation(str(OUT))
    trim_to_slides(prs, [IDX_COVER, IDX_CONTENT])

    fill_cover_slide(
        prs.slides[0],
        "[Presentation title]",
        "[Subtitle line 1]",
        "[Subtitle line 2 — owner / date]",
        "[Tag / workstream]",
    )
    fill_content_slide(
        prs.slides[1],
        "[Slide title]",
        [
            "[Bullet or paragraph 1]",
            "[Bullet or paragraph 2]",
            "[Bullet or paragraph 3]",
        ],
        subtitle="[Subtitle or lead line — optional]",
    )
    apply_cover_colors(prs.slides[0])
    apply_content_colors(prs.slides[1])

    save_presentation(prs, OUT)
    assert_pptx_valid(OUT)
    print(f"OK: {OUT} ({len(prs.slides)} slides, zip clean, no repair expected)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
