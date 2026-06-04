#!/usr/bin/env python3
"""Litmus checks for upscale-company-template.pptx — run before commit."""

import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "assets/templates/upscale-company-template.pptx"

sys.path.insert(0, str(ROOT / "scripts"))
from pptx_util import (  # noqa: E402
    COLOR_GOLD,
    COLOR_HEADLINE,
    COLOR_META,
    COLOR_NAVY_TITLE,
    COLOR_TAG,
    RIGHT_PANEL_LEFT,
    assert_pptx_valid,
    is_footer,
)


def _run_rgb(shape):
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                return run.font.color.rgb
            except AttributeError:
                continue
    return None


def main() -> int:
    errors = []
    if not TEMPLATE.exists():
        print(f"MISSING: {TEMPLATE}")
        return 1

    try:
        assert_pptx_valid(TEMPLATE)
    except ValueError as e:
        errors.append(str(e))

    prs = Presentation(str(TEMPLATE))
    if len(prs.slides) != 2:
        errors.append(f"expected 2 slides, got {len(prs.slides)}")

    cover = prs.slides[0]
    has_footer = any(
        "Upscale AI" in (sh.text_frame.text if sh.has_text_frame else "")
        for sh in cover.shapes
    )
    if not has_footer:
        errors.append("cover missing Upscale footer")

    right_gold = left_headline = meta_ok = tag_ok = False
    for sh in cover.shapes:
        if not sh.has_text_frame or is_footer(sh):
            continue
        t = sh.text_frame.text
        rgb = _run_rgb(sh)
        if sh.left >= RIGHT_PANEL_LEFT and "[Subtitle line 1]" in t:
            if rgb == COLOR_GOLD:
                right_gold = True
            else:
                errors.append(f"navy callout must be #FFE19E, got {rgb}")
        if "[Presentation title]" in t:
            if rgb == COLOR_HEADLINE:
                left_headline = True
            else:
                errors.append(f"cover title must be #FEDB8D gold, got {rgb}")
        if "[Subtitle line 2" in t:
            meta_ok = rgb == COLOR_META
        if "[Tag /" in t:
            tag_ok = rgb == COLOR_TAG

    if not right_gold:
        errors.append("cover: navy [Subtitle line 1] must be gold")
    if not left_headline:
        errors.append("cover: [Presentation title] must be headline gold #FEDB8D")
    if not meta_ok:
        errors.append("cover: meta line must be #AABBCC")
    if not tag_ok:
        errors.append("cover: tag line must be #778899")

    content = prs.slides[1]
    title_ok = body_ok = False
    for sh in content.shapes:
        if not sh.has_text_frame or is_footer(sh):
            continue
        if sh.top < 700_000 and "[Slide title]" in sh.text_frame.text:
            if _run_rgb(sh) == COLOR_NAVY_TITLE:
                title_ok = True
            else:
                errors.append(f"content title color wrong: {_run_rgb(sh)}")
        elif sh.top > 700_000 and "[Bullet" in sh.text_frame.text:
            body_ok = True
    if not title_ok:
        errors.append("content: [Slide title] must be navy #051830")
    if not body_ok:
        errors.append("content: missing bullet placeholders")

    if errors:
        print("VALIDATION FAILED:")
        for e in errors:
            print(f"  - {e}")
        return 1

    print(f"OK: {TEMPLATE}")
    print("  - ZIP clean, 2 slides, footer, gold on navy, navy title on content")
    return 0


if __name__ == "__main__":
    sys.exit(main())
