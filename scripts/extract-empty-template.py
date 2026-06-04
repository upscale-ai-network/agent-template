#!/usr/bin/env python3
"""
Strip all content slides from a source PPTX; keep slide masters/theme.
Writes Gluon empty template + template-spec.md.

Usage:
  python3 scripts/extract-empty-template.py /path/to/Mirror-Sflow-Bugatti-ASIC-CCC.pptx

Do not commit the source deck to agent-template unless IT/policy OK.
"""

import re
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))
from pptx_util import assert_pptx_valid, save_presentation  # noqa: E402
OUT_DIR = ROOT / "assets" / "templates"
OUT_PPTX = OUT_DIR / "upscale-exec-empty.pptx"
OUT_SPEC = OUT_DIR / "template-spec.md"

NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def delete_all_slides(prs: Presentation) -> int:
    sld_id_lst = prs.slides._sldIdLst
    count = len(sld_id_lst)
    for i in range(count - 1, -1, -1):
        r_id = sld_id_lst[i].rId
        prs.part.drop_rel(r_id)
        del sld_id_lst[i]
    return count


def _hex_from_color(el):
    if el is None:
        return None
    srgb = el.find("a:srgbClr", NS)
    if srgb is not None and "val" in srgb.attrib:
        return "#" + srgb.attrib["val"]
    scheme = el.find("a:schemeClr", NS)
    if scheme is not None and "val" in scheme.attrib:
        return f"theme:{scheme.attrib['val']}"
    return None


def extract_theme_hints(pptx_path: Path) -> dict:
    hints = {"slide_size": None, "theme_colors": [], "fonts": []}
    with zipfile.ZipFile(pptx_path) as z:
        for name in ("ppt/theme/theme1.xml", "ppt/presentation.xml"):
            if name not in z.namelist():
                continue
            root = ET.fromstring(z.read(name))
            if name.endswith("theme1.xml"):
                for clr in root.findall(".//a:clrScheme/a:*", NS):
                    tag = clr.tag.split("}")[-1]
                    hx = _hex_from_color(clr)
                    if hx:
                        hints["theme_colors"].append(f"{tag}: {hx}")
                for font in root.findall(".//a:fontScheme//*[@typeface]", NS):
                    tf = font.attrib.get("typeface")
                    if tf and tf not in hints["fonts"]:
                        hints["fonts"].append(tf)
            if name.endswith("presentation.xml"):
                sld_sz = root.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}sldSz")
                if sld_sz is not None:
                    cx, cy = sld_sz.get("cx"), sld_sz.get("cy")
                    hints["slide_size"] = f"{cx} x {cy} EMU"
    return hints


def write_spec(source: Path, removed: int, hints: dict) -> None:
    lines = [
        "# Template spec (auto-generated)",
        "",
        f"| Field | Value |",
        f"|-------|--------|",
        f"| Source deck | `{source.name}` (local only — not committed) |",
        f"| Source path | `{source}` |",
        f"| Content slides removed | {removed} |",
        f"| Git artifact | `upscale-exec-empty.pptx` |",
        f"| Slide size (EMU) | {hints.get('slide_size') or 'see PowerPoint'} |",
        "",
        "## Theme colors (from theme1.xml)",
        "",
    ]
    if hints["theme_colors"]:
        for c in hints["theme_colors"]:
            lines.append(f"- {c}")
    else:
        lines.append("- (none parsed — check Slide Master in PowerPoint)")
    lines += ["", "## Fonts (from theme1.xml)", ""]
    if hints["fonts"]:
        for f in hints["fonts"]:
            lines.append(f"- {f}")
    else:
        lines.append("- (none parsed)")
    lines += [
        "",
        "## Empty template checklist",
        "",
        "- [x] Content slides removed",
        "- [ ] Human: open `upscale-exec-empty.pptx` — confirm masters/layouts OK",
        "- [ ] Human: Save as `.potx` if you prefer template extension",
        "- [ ] Optional: export title slide thumbnail for README",
        "",
        "**Do not commit** the source Shafi/Bugatti deck to this remote unless approved.",
        "",
    ]
    OUT_SPEC.write_text("\n".join(lines), encoding="utf-8")


def main() -> int:
    if len(sys.argv) != 2:
        print("Usage: python3 scripts/extract-empty-template.py /path/to/source.pptx")
        return 1
    source = Path(sys.argv[1]).expanduser().resolve()
    if not source.is_file():
        print(f"Not found: {source}")
        return 1
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    hints = extract_theme_hints(source)
    prs = Presentation(str(source))
    removed = delete_all_slides(prs)
    if len(prs.slides) != 0:
        print("Warning: slides remain after delete")
    save_presentation(prs, OUT_PPTX)
    assert_pptx_valid(OUT_PPTX)
    write_spec(source, removed, hints)
    print(f"Removed {removed} slide(s) from {source.name}")
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_SPEC}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
