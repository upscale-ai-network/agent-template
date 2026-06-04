#!/usr/bin/env python3
"""Apply CCC-correct font colors to upscale-company-template.pptx in place."""

import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "assets/templates/upscale-company-template.pptx"

sys.path.insert(0, str(ROOT / "scripts"))
from pptx_util import apply_content_colors, apply_cover_colors, check_zip_duplicates  # noqa: E402


def main() -> int:
    prs = Presentation(str(TEMPLATE))
    apply_cover_colors(prs.slides[0])
    apply_content_colors(prs.slides[1])
    prs.save(str(TEMPLATE))
    dups = check_zip_duplicates(TEMPLATE)
    if dups:
        print("WARNING zip:", dups)
        return 1
    print(f"Colors fixed: {TEMPLATE}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
