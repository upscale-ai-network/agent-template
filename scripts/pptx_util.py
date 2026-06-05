"""Safe PPTX helpers — avoid duplicate layout parts; preserve Upscale text colors."""

from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Upscale CCC palette (from slide XML — do not use plain black on cover)
COLOR_GOLD = RGBColor(0xFF, 0xE1, 0x9E)  # logo, navy callout, footer logo
COLOR_HEADLINE = RGBColor(0xFE, 0xDB, 0x8D)  # cover main title (left)
COLOR_META = RGBColor(0xAA, 0xBB, 0xCC)  # cover owner/date line
COLOR_TAG = RGBColor(0x77, 0x88, 0x99)  # cover workstream tag
COLOR_NAVY_TITLE = RGBColor(0x05, 0x18, 0x30)  # content slide title
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)  # content body (not harsh #000 on white)
COLOR_LEAD = RGBColor(0x55, 0x66, 0x77)  # content optional subtitle line
COLOR_FOOTER = RGBColor(0x66, 0x88, 0xAA)  # © line on cover
COLOR_FOOTER_CONTENT = RGBColor(0x88, 0x88, 0x88)  # © line on content slides

RIGHT_PANEL_LEFT = 5_000_000  # EMU — shapes on navy background
FOOTER_TOP = 4_800_000
TITLE_TOP_MAX = 700_000

# Fonts — match CCC template (Calibri title, Arial body)
FONT_TITLE_FACE = "Calibri"
FONT_BODY_FACE = "Arial"
FONT_TITLE_SINGLE = Pt(26)
FONT_TITLE_MULTI = Pt(19)
FONT_LEAD = Pt(13)
FONT_BODY = Pt(13)
FONT_COVER_DSBM = Pt(22)

# Content slide layout — fixed bands on slides 1–4.
# Slide is 10in x 5.625in (16:9); title ~top 0.22, footer top ~5.34.
CONTENT_TITLE_BAND_BOTTOM = Inches(1.6)
CONTENT_DIAGRAM_TOP = Inches(1.65)
CONTENT_DIAGRAM_MAX_W = Inches(9.2)
CONTENT_DIAGRAM_MAX_H = Inches(3.55)


def delete_slide(prs: Presentation, index: int) -> None:
    """Remove slide at index from presentation."""
    slides = prs.slides._sldIdLst
    r_id = slides[index].rId
    prs.part.drop_rel(r_id)
    del slides[index]


def trim_to_slides(prs: Presentation, keep_indices: list) -> None:
    """Delete all slides whose index is not in keep_indices (0-based)."""
    keep = set(keep_indices)
    for i in range(len(prs.slides) - 1, -1, -1):
        if i not in keep:
            delete_slide(prs, i)


def is_footer(shape) -> bool:
    return shape.top > FOOTER_TOP


def _style_run(run, rgb: RGBColor, face: str, size: Pt) -> None:
    run.font.name = face
    run.font.size = size
    run.font.color.rgb = rgb


def _apply_run_color(paragraph, rgb: RGBColor) -> None:
    for run in paragraph.runs:
        run.font.color.rgb = rgb


def color_all_paragraphs(text_frame, rgb: RGBColor) -> None:
    for para in text_frame.paragraphs:
        _apply_run_color(para, rgb)


DSBM_TITLE_LINES = ["Dynamic", "Switch-Buffer", "Management"]


def set_shape_lines(
    shape,
    lines: List[str],
    rgb: RGBColor,
    center: bool = False,
    *,
    face: str = FONT_TITLE_FACE,
    size: Pt = FONT_TITLE_MULTI,
) -> None:
    """Replace text with multiple paragraphs (e.g. Switch-Buffer centered on line 2)."""
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(2)
        if center:
            p.alignment = PP_ALIGN.CENTER
        if line.strip():
            for run in p.runs:
                _style_run(run, rgb, face, size)


def _title_font_size(lines: List[str]) -> Pt:
    return FONT_TITLE_SINGLE if len(lines) == 1 else FONT_TITLE_MULTI


def set_content_title_block(
    shape,
    *,
    title: Optional[str] = None,
    title_lines: Optional[List[str]] = None,
    lead_line: Optional[str] = None,
) -> None:
    """Fixed title band: navy title + optional lead (slide 1)."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False  # keep title on one line (box is wider than text)
    entries: List[tuple[str, RGBColor, Pt, str]] = []
    if title_lines:
        for line in title_lines:
            entries.append((line, COLOR_NAVY_TITLE, FONT_TITLE_MULTI, FONT_TITLE_FACE))
    elif title:
        entries.append((title, COLOR_NAVY_TITLE, FONT_TITLE_SINGLE, FONT_TITLE_FACE))
    if lead_line:
        entries.append((lead_line, COLOR_LEAD, FONT_LEAD, FONT_BODY_FACE))
    for i, (text, rgb, size, face) in enumerate(entries):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(3)
        for run in p.runs:
            _style_run(run, rgb, face, size)


def set_shape_single_line(shape, text: str, rgb: Optional[RGBColor] = None) -> None:
    """Replace one line of text; apply explicit run color."""
    if rgb is None:
        rgb = COLOR_GOLD if shape.left >= RIGHT_PANEL_LEFT else COLOR_BODY
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
        for para in tf.paragraphs[1:]:
            for run in para.runs:
                run.text = ""
        _apply_run_color(tf.paragraphs[0], rgb)
    else:
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        _apply_run_color(p, rgb)


def fill_text_frame_lines(
    text_frame,
    lines: List[str],
    body_rgb: RGBColor,
    first_line_rgb: Optional[RGBColor] = None,
    *,
    lead_line_count: int = 0,
) -> None:
    text_frame.clear()
    lead_n = lead_line_count if lead_line_count else (1 if first_line_rgb else 0)
    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line
        p.level = 0
        if not line.strip():
            rgb = body_rgb
            size = FONT_BODY
        elif i < lead_n:
            rgb = first_line_rgb or COLOR_LEAD
            size = FONT_LEAD
        else:
            rgb = body_rgb
            size = FONT_BODY
        if line.strip():
            for run in p.runs:
                _style_run(run, rgb, FONT_BODY_FACE, size)


def apply_cover_colors(slide) -> None:
    """Fix colors on cover slide by shape position (placeholders or CCC text)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text
        if is_footer(shape):
            if shape.left < 5_000_000 and "Upscale AI" in t:
                color_all_paragraphs(shape.text_frame, COLOR_FOOTER)
            elif "UP" in t or "upscale" in t:
                color_all_paragraphs(shape.text_frame, COLOR_GOLD)
            continue
        if shape.left >= RIGHT_PANEL_LEFT and shape.top < 2_500_000:
            color_all_paragraphs(shape.text_frame, COLOR_GOLD)
        elif shape.top > 3_000_000 and shape.top < 3_600_000:
            color_all_paragraphs(shape.text_frame, COLOR_META)
        elif shape.top > 3_600_000 and shape.top < 4_200_000:
            color_all_paragraphs(shape.text_frame, COLOR_TAG)
        elif shape.top > 1_500_000 and shape.top < 2_500_000 and shape.left < RIGHT_PANEL_LEFT:
            color_all_paragraphs(shape.text_frame, COLOR_HEADLINE)
        elif "UP" in t and shape.top < 1_500_000:
            color_all_paragraphs(shape.text_frame, COLOR_GOLD)


def apply_content_colors(slide, *, skip_title_band: bool = False) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if skip_title_band and shape.top < TITLE_TOP_MAX:
            continue
        if is_footer(shape):
            if shape.left < 5_000_000 and "Upscale AI" in shape.text_frame.text:
                color_all_paragraphs(shape.text_frame, COLOR_FOOTER_CONTENT)
            elif "UP" in shape.text_frame.text or "upscale" in shape.text_frame.text:
                color_all_paragraphs(shape.text_frame, COLOR_GOLD)
            continue
        if shape.top < TITLE_TOP_MAX:
            color_all_paragraphs(shape.text_frame, COLOR_NAVY_TITLE)
        else:
            tf = shape.text_frame
            for i, para in enumerate(tf.paragraphs):
                line = para.text.strip()
                if not line:
                    continue
                if i == 0 and ("optional" in line or "Subtitle" in line):
                    _apply_run_color(para, COLOR_LEAD)
                else:
                    _apply_run_color(para, COLOR_BODY)


def fill_cover_slide(slide, headline: str, subline: str, meta: str, tag: str) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame or is_footer(shape):
            continue
        raw = shape.text_frame.text
        if "Mirror CCC" in raw or "[Presentation title]" in raw:
            set_shape_single_line(shape, headline, COLOR_HEADLINE)
        elif "Bugatti" in raw or "[Subtitle line 1]" in raw:
            if "\n" in subline:
                set_shape_lines(shape, subline.split("\n"), COLOR_GOLD, center=True)
            elif subline == " ".join(DSBM_TITLE_LINES):
                set_shape_lines(
                    shape,
                    DSBM_TITLE_LINES,
                    COLOR_GOLD,
                    center=True,
                    face=FONT_TITLE_FACE,
                    size=FONT_COVER_DSBM,
                )
            else:
                set_shape_lines(shape, [subline], COLOR_GOLD, center=True)
        elif raw.strip().startswith("Authors:") or "[Subtitle line 2" in raw:
            set_shape_single_line(shape, meta, COLOR_META)
        elif "Hardware/Architecture" in raw or "[Tag /" in raw:
            set_shape_single_line(shape, tag, COLOR_TAG)
    apply_cover_colors(slide)


def fill_content_slide(
    slide,
    title: str,
    bullets: List[str],
    subtitle: Optional[str] = None,
    lead: Optional[str] = None,
    title_lines: Optional[List[str]] = None,
) -> None:
    title_shape = None
    body_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame or is_footer(shape):
            continue
        if shape.top < 700_000:
            title_shape = shape
        else:
            body_shape = shape
    if title_shape:
        if title_lines:
            set_shape_lines(
                title_shape,
                title_lines,
                COLOR_NAVY_TITLE,
                center=False,
                size=_title_font_size(title_lines),
            )
        elif "\n" in title:
            parts = title.split("\n")
            set_shape_lines(
                title_shape,
                parts,
                COLOR_NAVY_TITLE,
                size=_title_font_size(parts),
            )
        else:
            set_shape_single_line(title_shape, title, COLOR_NAVY_TITLE)
            for run in title_shape.text_frame.paragraphs[0].runs:
                _style_run(run, COLOR_NAVY_TITLE, FONT_TITLE_FACE, FONT_TITLE_SINGLE)
    if body_shape:
        lines: List[str] = []
        lead_count = 0
        if subtitle:
            lines.append(subtitle)
            lead_count += 1
        if lead:
            lines.append(lead)
            lead_count += 1
        if lead_count:
            lines.append("")
        lines.extend(bullets)
        fill_text_frame_lines(
            body_shape.text_frame,
            lines,
            COLOR_BODY,
            first_line_rgb=COLOR_LEAD if lead_count else None,
            lead_line_count=lead_count,
        )
    apply_content_colors(slide)


def hide_body_placeholder(slide) -> None:
    """Hide bullet placeholder so only diagram shows in content area."""
    for shape in slide.shapes:
        if not shape.has_text_frame or is_footer(shape):
            continue
        if shape.top >= TITLE_TOP_MAX:
            shape.text_frame.clear()
            shape.width = 0
            shape.height = 0


def place_diagram(slide, image_path: Path, slide_width: Optional[int] = None) -> None:
    """Fit diagram in fixed box — same visual size on every slide."""
    slide_w = slide_width or int(Inches(13.333))
    box_top = int(CONTENT_DIAGRAM_TOP)
    box_h = int(CONTENT_DIAGRAM_MAX_H)
    max_w = int(CONTENT_DIAGRAM_MAX_W)
    pic = slide.shapes.add_picture(str(image_path), 0, box_top)
    scale_w = max_w / pic.width
    scale_h = box_h / pic.height
    scale = scale_w if pic.height * scale_w <= box_h else min(scale_w, scale_h)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int((slide_w - pic.width) / 2)
    pic.top = box_top + (box_h - pic.height) // 2


def fill_content_diagram_slide(
    slide,
    title: str,
    image_path: Path,
    *,
    title_lines: Optional[List[str]] = None,
    lead_line: Optional[str] = None,
    slide_width: Optional[int] = None,
) -> None:
    """Title band + optional lead + diagram in aligned box; no bullet body."""
    title_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame or is_footer(shape):
            continue
        if shape.top < TITLE_TOP_MAX:
            title_shape = shape
            break
    if title_shape:
        set_content_title_block(
            title_shape,
            title=title,
            title_lines=title_lines,
            lead_line=lead_line,
        )
    hide_body_placeholder(slide)
    if not image_path.is_file():
        raise FileNotFoundError(f"Diagram slide requires image: {image_path}")
    place_diagram(slide, image_path, slide_width=slide_width)
    apply_content_colors(slide, skip_title_band=True)
    if lead_line and title_shape and title_shape.text_frame.paragraphs:
        p = title_shape.text_frame.paragraphs[-1]
        for run in p.runs:
            _style_run(run, COLOR_LEAD, FONT_BODY_FACE, FONT_LEAD)


def check_zip_duplicates(path) -> list:
    """Return duplicate member names in pptx zip (empty = OK)."""
    import zipfile
    from collections import Counter

    path = Path(path)
    if path.name.startswith("~$") or not path.is_file():
        return []
    with zipfile.ZipFile(path) as z:
        counts = Counter(z.namelist())
        return [n for n, c in counts.items() if c > 1]


def assert_pptx_valid(path) -> None:
    """
    Raise ValueError if PowerPoint is likely to show Repair.
    Checks: zip integrity, no duplicate parts, python-pptx can open.
    """
    import zipfile

    path = Path(path)
    if path.name.startswith("~$"):
        raise ValueError(f"refusing PowerPoint lock file: {path}")
    if not path.is_file() or path.stat().st_size < 1024:
        raise ValueError(f"not a valid pptx file: {path}")

    with zipfile.ZipFile(path) as z:
        corrupt = z.testzip()
        if corrupt:
            raise ValueError(f"zip corrupt member {corrupt!r} in {path.name}")

    dups = check_zip_duplicates(path)
    if dups:
        raise ValueError(f"duplicate zip entries in {path.name}: {dups}")

    Presentation(str(path))


def save_presentation(prs: Presentation, path) -> Path:
    """Atomic save + assert before replacing target (no corrupt file on disk)."""
    path = Path(path)
    tmp = path.with_suffix(".pptx.tmp")
    try:
        prs.save(str(tmp))
        assert_pptx_valid(tmp)
        tmp.replace(path)
    finally:
        if tmp.exists() and not path.exists():
            tmp.replace(path)
        elif tmp.exists():
            tmp.unlink()
    return path
