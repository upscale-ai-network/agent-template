#!/usr/bin/env python3
"""Dump on-slide text per slide — cold-read aid for md→pptx decks."""

from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from pptx import Presentation  # noqa: E402


def dump(path: Path) -> str:
    prs = Presentation(str(path))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"=== Slide {i} ===")
        chunks: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                chunks.append(text)
        if chunks:
            lines.extend(chunks)
        else:
            lines.append("(no text shapes — diagram-only)")
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def main() -> int:
    if len(sys.argv) < 2:
        print("Usage: dump_deck_on_slide.py path/to/deck.pptx", file=sys.stderr)
        return 1
    path = Path(sys.argv[1])
    if not path.is_file():
        print(f"Not found: {path}", file=sys.stderr)
        return 1
    sys.stdout.write(dump(path))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
