"""Contra-tests for pptx_util — P1 guard, not a spec appendix.

Forensics note (2026-06-19, Gluon / diwakar-work lattice session):

TRIAGE FRAMEWORK (agreed that day)
  Three buckets for test failures:
    #1 Production — confirmed prod bugs; must fix (finite; subset of #3).
    #2 Test gaps — missing coverage, infra, docs; infinite by definition.
    #3 Unknown — finite list of current fails; triage each into #1, #2, or dismiss.
  Relation: #1 ⊆ #3. P1 lived in #3 until reproduced, then promoted to #1.

BUG P1 (promoted to #1)
  File: scripts/pptx_util.py::fill_content_slide
  deck_from_md emits BulletItem = (level, text) tuples. Diagram/split slides
  normalize via _normalize_bullets in fill_content_diagram_slide; text-only
  slides call lines.extend(bullets) then fill_text_frame_lines(List[str]).
  Symptom: TypeError: expected string or bytes-like object, got 'tuple'
  at p.text = line (python-pptx / re.split on non-string).
  Trigger: workflow canary slide 3 (text-only + bullets). A3/B6 decks do not
  hit this path today (diagram/image only). Workflow tests tc04/tc12 masked
  until Puppeteer Chrome available — they die at mermaid before slide 3.

PROCESS (explicit owner directive: do NOT fix prod first)
  1. Write contra-test stating contract (parsed bullets render on text-only slides).
  2. Run it; MUST fail on current prod — falsify SoT, prove bug is real.
  3. Only then patch prod (likely: _normalize_bullets in fill_content_slide).
  4. Same test goes green; keep permanently as validator (see below).

NOT A FALSE FAILURE
  "False failure" = bad test, flaky env, wrong setup.
  This test is correct; prod is wrong. Intentional red until P1 patch.
  Default suite while open: was 30 tests / 8 fails → 31 tests / 9 fails (+this).

TWO VALID FALSIFICATION STYLES (discussed, one chosen)
  A) Owner's framing: contra-test as temporary probe; fix prod; prune test so
     suite stays derived from first-principles specs, not evolutionary appendages.
     More epistemically pure (epistēmē = knowledge; probe = how we know).
  B) Agent's framing: keep contra-test forever as named contract/regression guard.
  Resolution (owner, same session): keep the test — "absolute purity can be
  poisonous; perfect enemy of good." Pruning is its own failure mode (bug
  returns, re-lands in cat-2/3 chaos). Named tripping point = documented
  warning for anyone willing to repeat history ~2 years out. Pragmatic blend
  noted: could later add spec-level canary slide-3 test and narrow this one,
  but do not delete on P1 fix.

PROPOSED FIX (NOT APPLIED YET — recorded for when contra-test is green)
  ~5–10 lines in fill_content_slide: route bullets through _normalize_bullets,
  format as bullet strings (match split-slide prefix rules). Low risk to A3/B6
  (path unused); high reward (unblocks text-only md decks + workflow canary).

OTHER FAILURES THAT DAY (not this file — context only)
  #3 triaged: 7× workflow/deck fails = Chrome/Puppeteer env; 1× tc16[b6] =
  stale published-deck-hashes.txt (regen ≡ committed, both ≠ golden). None
  confirmed prod besides P1.
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest
from pptx import Presentation

from conftest import p1_verify  # noqa: E402

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from pptx_util import fill_content_slide  # noqa: E402

STYLE_REF = ROOT / "assets/templates/upscale-ccc-style-reference.pptx"


@p1_verify
@pytest.mark.skipif(not STYLE_REF.is_file(), reason="company style template missing")
def test_fill_content_slide_accepts_deck_from_md_bullet_tuples():
    """P1 contra-test: tuples from deck_from_md must not crash text-only slides."""
    prs = Presentation(str(STYLE_REF))
    slide = prs.slides[2]
    # Same shape as deck_from_md._parse_bullet — not raw markdown strings.
    bullets = [(0, "Top bullet"), (1, "Sub bullet")]
    fill_content_slide(
        slide,
        "Text-only slide",
        bullets,
        subtitle="Act subtitle",
        lead="Lead line",
    )
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    body = "\n".join(t for t in texts if "Top bullet" in t or "Sub bullet" in t)
    assert "Top bullet" in body
    assert "Sub bullet" in body
